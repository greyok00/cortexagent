#!/usr/bin/env python3
"""lib/document_adapter.py — extract text from documents (PDF/DOCX/XLSX/PPTX).

Docling if installed (tables, reading order, OCR for scans), else format
fallbacks: pdftotext for PDFs, python-docx for DOCX, openpyxl for XLSX,
python-pptx for PPTX. Every path degrades gracefully — the adapter never
raises; it returns {"ok", "text", "error"}.

Usage:
  python3 lib/document_adapter.py --smoke
"""
from __future__ import annotations

import subprocess
import sys
from pathlib import Path
from typing import Any, Dict

_REPO_ROOT = Path(__file__).resolve().parent.parent
if str(_REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(_REPO_ROOT))


def _extract_pdf(file: str) -> str:
    """PDF text via Docling (if installed) else pdftotext -layout."""
    try:
        from docling.document_converter import DocumentConverter
        result = DocumentConverter().convert(file)
        return result.document.export_to_text()
    except Exception:
        pass  # Docling missing or failed — fall back to pdftotext
    try:
        out = subprocess.run(["pdftotext", "-layout", file, "-"],
                             capture_output=True, text=True, timeout=120)
        return out.stdout
    except Exception:
        return ""


def _extract_docx(file: str) -> str:
    try:
        import docx
        return "\n".join(p.text for p in docx.Document(file).paragraphs)
    except Exception:
        return ""


def _extract_xlsx(file: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file, read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append("\t".join(cells))
        return "\n".join(rows)
    except Exception:
        return ""


def _extract_pptx(file: str) -> str:
    try:
        from pptx import Presentation
        parts = []
        for slide in Presentation(file).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    except Exception:
        return ""


def parse_document(file: str) -> Dict[str, Any]:
    """Extract text from a document. Returns {"ok", "text", "error"}."""
    path = Path(file)
    if not path.is_file():
        return {"ok": False, "text": "", "error": f"file not found: {file}"}
    ext = path.suffix.lower()
    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
    }
    fn = extractors.get(ext)
    if fn is None:
        return {"ok": False, "text": "", "error": f"unsupported format: {ext}"}
    text = fn(str(path)).strip()
    if not text:
        return {"ok": True, "text": "",
                "error": "no text layer (scanned PDF? OCR unavailable)"}
    return {"ok": True, "text": text, "error": ""}


def _smoke() -> int:
    fails = 0
    import shutil
    import tempfile
    from pathlib import Path
    tmp = Path(tempfile.mkdtemp())
    try:
        # sample PDF with a text layer via reportlab
        from reportlab.pdfgen import canvas
        pdf = tmp / "sample.pdf"
        c = canvas.Canvas(str(pdf))
        c.drawString(72, 720, "CortexAgent document adapter smoke test")
        c.save()
        r = parse_document(str(pdf))
        if not r.get("ok") or "smoke test" not in r.get("text", ""):
            print(f"❌ pdf parse: {r}")
            fails += 1
        # missing file → clean error
        r = parse_document(str(tmp / "nope.pdf"))
        if r.get("ok") or "not found" not in r.get("error", ""):
            print(f"❌ missing file: {r}")
            fails += 1
        # unsupported format → clean error
        (tmp / "x.xyz").write_text("junk")
        r = parse_document(str(tmp / "x.xyz"))
        if r.get("ok") or "unsupported" not in r.get("error", ""):
            print(f"❌ unsupported: {r}")
            fails += 1
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
    print("document_adapter smoke PASS" if fails == 0 else f"❌ {fails} failures")
    return 1 if fails else 0


def main() -> int:
    if len(sys.argv) > 1 and sys.argv[1] == "--smoke":
        return _smoke()
    print("Usage: python3 lib/document_adapter.py --smoke")
    return 0


if __name__ == "__main__":
    sys.exit(main())
